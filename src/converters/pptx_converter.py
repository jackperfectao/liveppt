"""PowerPoint (.pptx) to LivePPT Markdown converter."""

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.converters.base import BaseConverter

logger = logging.getLogger(__name__)


class PptxConverter(BaseConverter):
    """Convert PowerPoint .pptx presentations to LivePPT-compatible Markdown.

    Mapping:
    - First title slide -> # (deck title) + subtitle as intro paragraph
    - Slide titles -> ## (slide boundary)
    - Body text shapes -> - bullets
    - Speaker notes -> > blockquote comment
    """

    def convert(self, input_path: Path) -> str:
        try:
            prs = Presentation(str(input_path))
        except Exception as e:
            raise ValueError(f"Cannot open .pptx file: {input_path}: {e}")

        lines: list[str] = []
        title_slide_done = False
        slide_number = 0

        for slide in prs.slides:
            slide_number += 1

            # Try to identify the slide title
            title_text = self._find_slide_title(slide)

            if not title_text:
                # Fallback: use slide number
                title_text = f"第 {slide_number} 页"

            # First slide special handling
            if not title_slide_done:
                layout_name = slide.slide_layout.name.lower() if slide.slide_layout else ""
                if "title" in layout_name or slide_number == 1:
                    title_slide_done = True
                    lines.append(f"# {title_text}")
                    lines.append("")

                    # Extract subtitle from second text part
                    subtitle = self._find_subtitle(slide)
                    if subtitle:
                        lines.append(subtitle)
                        lines.append("")
                    continue

            # Regular content slide
            lines.append(f"## {title_text}")
            lines.append("")

            # Extract bullets from remaining text
            bullets = self._extract_bullets(slide, title_text)
            for bullet in bullets:
                lines.append(f"- {bullet}")
            if bullets:
                lines.append("")

            # Speaker notes
            notes = self._extract_notes(slide)
            if notes:
                lines.append(f"> {notes}")
                lines.append("")

        if not lines:
            return self._empty_deck(str(input_path.stem))

        return "\n".join(lines).strip() + "\n"

    def _find_slide_title(self, slide) -> Optional[str]:
        """Try to find the slide title from placeholder shapes or first text box."""
        for shape in slide.shapes:
            if shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                # title or center_title
                if placeholder_type is not None and str(placeholder_type) in ("1", "3"):
                    text = shape.text_frame.text.strip()
                    if text:
                        return text

        # Fallback: first non-empty text shape that's reasonably short
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) < 200:
                    return text
        return None

    def _find_subtitle(self, slide) -> Optional[str]:
        """Extract subtitle text from a title slide."""
        for shape in slide.shapes:
            if shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                if placeholder_type is not None and str(placeholder_type) == "2":  # subtitle
                    text = shape.text_frame.text.strip()
                    if text:
                        return text
        return None

    def _extract_bullets(self, slide, title_text: str) -> list[str]:
        """Extract body text from a slide as bullet items."""
        bullets = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                if placeholder_type is not None and str(placeholder_type) in ("1", "3"):
                    continue  # skip title

            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text != title_text:
                    # Split multi-paragraph text frames into separate bullets
                    for para in shape.text_frame.paragraphs:
                        pt = para.text.strip()
                        if pt and pt != title_text:
                            bullets.append(pt)

            if shape.shape_type == MSO_SHAPE_TYPE.TABLE and shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        bullets.append(" | ".join(cells))

        return bullets

    def _extract_notes(self, slide) -> Optional[str]:
        """Extract speaker notes if available."""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    text = notes_slide.notes_text_frame.text.strip()
                    if text:
                        return text
        except Exception:
            pass
        return None

    def _empty_deck(self, title: str) -> str:
        return f"# {title}\n\n演示文稿内容为空或无法解析。\n"
